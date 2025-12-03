from datetime import datetime
from pathlib import Path
from typing import Iterable, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from charts import ChartArtifact
from insights import InsightSections

OUTPUT_PATH = Path("data") / "output" / "InsightForge_Report.pptx"

BODY_FONT = "Calibri"
ACCENT_RGB = RGBColor(31, 78, 121)


def _add_bullets(placeholder, lines: List[str]) -> None:
    placeholder.text = ""
    tf = placeholder.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        if idx == 0:
            tf.text = line
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.color.rgb = ACCENT_RGB
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = ACCENT_RGB


def _add_cover_slide(prs: Presentation, timestamp: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "InsightForge Executive Brief"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ACCENT_RGB
    subtitle = slide.placeholders[1]
    subtitle.text = f"Automated insight pack • {timestamp}"


def _add_section_slide(prs: Presentation, title: str, bullets: List[str], limit: int = 6) -> bool:
    clean = [b for b in bullets if b.strip()]
    if not clean:
        return False
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ACCENT_RGB
    _add_bullets(slide.shapes.placeholders[1], clean[:limit])
    return True


def _add_chart_slide(prs: Presentation, chart: ChartArtifact) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = chart.title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ACCENT_RGB
    left = Inches(0.9)
    top = Inches(1.2)
    slide.shapes.add_picture(str(chart.path), left, top, width=Inches(8))

    tx_box = slide.shapes.add_textbox(Inches(0.9), Inches(6.2), Inches(8), Inches(1))
    tf = tx_box.text_frame
    tf.text = chart.description
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = ACCENT_RGB


def build_presentation(insights: InsightSections, charts: Iterable[ChartArtifact]) -> Path:
    prs = Presentation()
    timestamp = datetime.now().strftime("%d %b %Y %H:%M")

    _add_cover_slide(prs, timestamp)
    _add_section_slide(prs, "Executive Overview", insights.overview)
    _add_section_slide(prs, "Key Metrics", insights.key_metrics)
    _add_section_slide(prs, "Trend Signals", insights.trends)
    _add_section_slide(prs, "Anomalies & Watchouts", insights.anomalies, limit=4)
    _add_section_slide(prs, "Strategic Recommendations", insights.recommendations, limit=4)

    for chart in charts:
        _add_chart_slide(prs, chart)

    _add_section_slide(prs, "Closing Summary", insights.summary)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH